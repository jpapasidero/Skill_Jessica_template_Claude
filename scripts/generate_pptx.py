#!/usr/bin/env python3
"""
Générateur de présentation au template "Palette Jessica" (Safran).

Usage :
    python generate_pptx.py --content content.json --output out.pptx [--logo logo.png]

Le fichier content.json contient une liste de slides, chacune décrivant son
layout (1..11) et le contenu textuel à injecter dans les placeholders nommés.

Voir references/content_schema.md pour le contrat d'entrée détaillé.
"""
import argparse
import json
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Cm

# Permettre l'exécution directe ET en module
try:
    from .builders import build_slide
    from .impact_parser import generate_impact_content
except ImportError:
    sys.path.insert(0, str(Path(__file__).parent))
    from builders import build_slide
    from impact_parser import generate_impact_content


# Format slide : 33.87 × 19.05 cm (16:9 widescreen)
SLIDE_W_CM = 33.867
SLIDE_H_CM = 19.05
PROJECT_ROOT = Path(__file__).resolve().parent.parent
LOGO_EXTENSIONS = (".png", ".jpg", ".jpeg")


def init_presentation():
    """Crée une présentation vierge au bon format."""
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W_CM)
    prs.slide_height = Cm(SLIDE_H_CM)
    return prs


def resolve_logo_path(logo_path=None):
    """Trouve le logo depuis un chemin explicite ou les dossiers d'assets locaux."""
    if logo_path:
        candidate = Path(logo_path).expanduser()
        search_paths = [candidate]
        if not candidate.is_absolute():
            search_paths.append(PROJECT_ROOT / candidate)
        for path in search_paths:
            if path.exists():
                return str(path)
        print(f"[WARN] Logo introuvable : {logo_path}", file=sys.stderr)

    logo_dirs = [
        PROJECT_ROOT / "assets" / "logo",
        PROJECT_ROOT / "assets" / "logos",
        PROJECT_ROOT / "asset" / "logo",
        PROJECT_ROOT / "asset" / "logos",
    ]
    preferred_names = [
        "safran_logo.png",
        "logo.png",
        "safran.png",
    ]

    for logo_dir in logo_dirs:
        if logo_dir.is_file() and logo_dir.suffix.lower() in LOGO_EXTENSIONS:
            return str(logo_dir)
        if not logo_dir.is_dir():
            continue
        for name in preferred_names:
            candidate = logo_dir / name
            if candidate.exists():
                return str(candidate)
        for candidate in sorted(logo_dir.iterdir()):
            if candidate.is_file() and candidate.suffix.lower() in LOGO_EXTENSIONS:
                return str(candidate)
    return None


def generate(content_path, output_path, logo_path=None, impact_path=None):
    logo_path = resolve_logo_path(logo_path)

    with open(content_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    slides_data = data.get("slides", [])

    # Si un fichier d'analyse d'impact est fourni, générer les slides Layout 15
    if impact_path:
        impact_slides = generate_impact_content(impact_path)
        if impact_slides:
            # Retirer les éventuels Layout 15 existants du contenu
            slides_data = [sd for sd in slides_data if sd.get("layout") != 15]
            slides_data.extend(impact_slides)
            print(f"[INFO] {len(impact_slides)} slide(s) Layout 15 générée(s) "
                  f"depuis {impact_path}", file=sys.stderr)

    prs = init_presentation()
    total = len(slides_data)

    # layout vide ("blank") pour chaque slide
    blank_layout = prs.slide_layouts[6]  # généralement "Blank" en index 6

    for i, sd in enumerate(slides_data, start=1):
        layout_n = sd.get("layout")
        if layout_n is None:
            print(f"[WARN] slide {i} : champ 'layout' manquant, skip", file=sys.stderr)
            continue
        slide = prs.slides.add_slide(blank_layout)
        # nettoyer les placeholders éventuels du layout blank (pour partir vraiment de zéro)
        for shp in list(slide.placeholders):
            sp = shp._element
            sp.getparent().remove(sp)

        try:
            build_slide(layout_n, slide, sd.get("content", {}),
                        page_num=i, total=total, logo_path=logo_path)
        except Exception as e:
            print(f"[ERROR] slide {i} layout {layout_n} : {e}", file=sys.stderr)
            raise

    prs.save(output_path)
    print(f"[OK] Présentation générée : {output_path} ({total} slides)")


def main():
    ap = argparse.ArgumentParser(description="Générateur PPTX Palette Jessica")
    ap.add_argument("--content", required=True, help="Fichier JSON de contenu")
    ap.add_argument("--output", required=True, help="Fichier .pptx de sortie")
    ap.add_argument("--logo", default=None,
                    help="Chemin optionnel vers le logo Safran. Par defaut, cherche dans assets/logo(s).")
    args = ap.parse_args()
    generate(args.content, args.output, logo_path=args.logo)


if __name__ == "__main__":
    main()
