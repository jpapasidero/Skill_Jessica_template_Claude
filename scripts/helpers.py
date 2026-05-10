"""
Helpers bas-niveau pour le template Palette Jessica.
Manipule directement le XML OOXML pour les effets non gérés nativement par python-pptx :
- cap="small" (petites majuscules)
- alpha (filigrane translucide)
- gradient vertical 3-stops sur le bandeau Take-away
"""
import re
from io import BytesIO
from pathlib import Path

from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image, ImageDraw, ImageFont


# ---------- Conversions ----------

def cm(v):
    return Cm(v)


def hex_to_rgb(hexstr):
    h = hexstr.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def hex_to_rgba(hexstr, alpha_pct=100):
    h = hexstr.lstrip("#")
    alpha = max(0, min(255, round(255 * alpha_pct / 100)))
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16), alpha)


def remove_theme_style(shape):
    """Remove PowerPoint theme effects that can add shadows by default."""
    style = shape._element.find(qn("p:style"))
    if style is not None:
        shape._element.remove(style)


# ---------- Texte ----------

def set_run_props(run, *, font=None, size_pt=None, bold=None, italic=None, color_hex=None,
                  cap=None, alpha_pct=None):
    """Applique les propriétés à un run et patche le XML pour cap / alpha si demandé."""
    if font:
        run.font.name = font
    if size_pt is not None:
        run.font.size = Pt(size_pt)
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if color_hex:
        run.font.color.rgb = hex_to_rgb(color_hex)

    rPr = run._r.get_or_add_rPr()

    # cap="small" : petites majuscules avec 1ère lettre majuscule (rendu OOXML natif)
    if cap == "small":
        rPr.set("cap", "small")
    elif cap == "all":
        rPr.set("cap", "all")

    # alpha : injecter <a:alpha val="N"/> dans <a:solidFill><a:srgbClr>
    if alpha_pct is not None and color_hex:
        # on retire le solidFill existant pour le reconstruire avec alpha
        for sf in rPr.findall(qn("a:solidFill")):
            rPr.remove(sf)
        sf = etree.SubElement(rPr, qn("a:solidFill"))
        srgb = etree.SubElement(sf, qn("a:srgbClr"))
        srgb.set("val", color_hex.lstrip("#").upper())
        alpha = etree.SubElement(srgb, qn("a:alpha"))
        alpha.set("val", str(int(alpha_pct * 1000)))  # 40% -> 40000


def normalize_text_spec(text):
    """Accepte une chaine simple ou une spec enrichie de type {"text": ..., ...}."""
    if isinstance(text, dict):
        spec = dict(text)
        spec["text"] = str(spec.get("text", ""))
        spec["emphasis"] = spec.get("emphasis", spec.get("emphases", [])) or []
        return spec
    return {"text": "" if text is None else str(text), "emphasis": []}


def normalize_line_spacing(value):
    """Convertit l'interligne JSON en valeur python-pptx (ratio ou points)."""
    if value is None:
        return None
    if isinstance(value, str):
        cleaned = value.strip()
        if cleaned.endswith("%"):
            return float(cleaned[:-1]) / 100.0
        if cleaned.endswith("pt"):
            return Pt(float(cleaned[:-2]))
        return float(cleaned)
    return value


def emphasis_spans(text, emphases):
    """Retourne des segments (texte, style) en appliquant les emphases non chevauchantes."""
    spans = []
    for emph in emphases:
        if not isinstance(emph, dict):
            emph = {"text": str(emph)}
        needle = str(emph.get("text", ""))
        if not needle:
            continue
        flags = 0 if emph.get("case_sensitive", False) else re.IGNORECASE
        count = 1 if emph.get("first_only", False) else 0
        for match in re.finditer(re.escape(needle), text, flags):
            start, end = match.span()
            if any(start < existing_end and end > existing_start for existing_start, existing_end, _ in spans):
                continue
            spans.append((start, end, emph))
            if count == 1:
                break

    if not spans:
        return [(text, None)]

    segments = []
    cursor = 0
    for start, end, emph in sorted(spans, key=lambda item: item[0]):
        if start > cursor:
            segments.append((text[cursor:start], None))
        segments.append((text[start:end], emph))
        cursor = end
    if cursor < len(text):
        segments.append((text[cursor:], None))
    return segments


def run_style(base, override=None, defaults=None):
    """Fusionne le style du placeholder avec une emphase locale."""
    style = dict(base)
    if not override:
        return style
    if defaults:
        for key, value in defaults.items():
            if key == "color":
                style["color_hex"] = value
            else:
                style[key] = value
    mapping = {
        "font": "font",
        "size_pt": "size_pt",
        "bold": "bold",
        "italic": "italic",
        "color": "color_hex",
        "color_hex": "color_hex",
        "cap": "cap",
        "alpha_pct": "alpha_pct",
    }
    for src, dest in mapping.items():
        if src in override:
            style[dest] = override[src]
    return style


def add_textbox(slide, x_cm, y_cm, w_cm, h_cm, text, *, font="Segoe UI", size_pt=12,
                bold=False, italic=False, color_hex="#070E1D", cap=None, alpha_pct=None,
                align="left", anchor="t", name=None, line_spacing=None,
                emphasis_style=None, margin_left=None, margin_right=None,
                margin_top=None, margin_bottom=None):
    """Crée un textbox avec texte simple ou spec enrichie (interligne + emphases)."""
    tb = slide.shapes.add_textbox(cm(x_cm), cm(y_cm), cm(w_cm), cm(h_cm))
    if name:
        tb.name = name
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = cm(margin_left) if margin_left is not None else Emu(45720)
    tf.margin_right = cm(margin_right) if margin_right is not None else Emu(45720)
    tf.margin_top = cm(margin_top) if margin_top is not None else Emu(45720)
    tf.margin_bottom = cm(margin_bottom) if margin_bottom is not None else Emu(45720)

    # ancrage vertical
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    anchor_map = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}
    tf.vertical_anchor = anchor_map.get(anchor, MSO_ANCHOR.TOP)

    p = tf.paragraphs[0]
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    p.alignment = align_map.get(align, PP_ALIGN.LEFT)

    spec = normalize_text_spec(text)
    effective_line_spacing = spec.get("line_spacing", spec.get("interligne", line_spacing))
    effective_line_spacing = normalize_line_spacing(effective_line_spacing)
    if effective_line_spacing is not None:
        p.line_spacing = effective_line_spacing

    base_style = {
        "font": font,
        "size_pt": size_pt,
        "bold": bold,
        "italic": italic,
        "color_hex": color_hex,
        "cap": cap,
        "alpha_pct": alpha_pct,
    }
    for segment, emphasis in emphasis_spans(spec["text"], spec["emphasis"]):
        if not segment:
            continue
        run = p.add_run()
        run.text = segment
        set_run_props(run, **run_style(base_style, emphasis, emphasis_style))
    return tb


# ---------- Formes ----------

def add_rect(slide, x_cm, y_cm, w_cm, h_cm, fill_hex, *, line_hex=None, line_w_pt=0,
             rounded=False, radius_adjust=0.1, name=None, text=None,
             font="Segoe UI", size_pt=12, bold=False, italic=False,
             text_color_hex="#070E1D", color_hex=None, cap=None, alpha_pct=None,
             align="left", anchor="m", line_spacing=None, emphasis_style=None,
             margin_left=None, margin_right=None, margin_top=None, margin_bottom=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, cm(x_cm), cm(y_cm), cm(w_cm), cm(h_cm))
    remove_theme_style(shp)
    if name:
        shp.name = name
    if rounded:
        # arrondi modéré
        shp.adjustments[0] = radius_adjust
    shp.fill.solid()
    shp.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    if line_hex:
        shp.line.color.rgb = hex_to_rgb(line_hex)
        shp.line.width = Pt(line_w_pt) if line_w_pt else Pt(0.25)
    else:
        shp.line.fill.background()
    if text is not None:
        tf = shp.text_frame
        tf.word_wrap = True
        tf.margin_left = cm(margin_left) if margin_left is not None else Emu(45720)
        tf.margin_right = cm(margin_right) if margin_right is not None else Emu(45720)
        tf.margin_top = cm(margin_top) if margin_top is not None else Emu(45720)
        tf.margin_bottom = cm(margin_bottom) if margin_bottom is not None else Emu(45720)

        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        anchor_map = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}
        tf.vertical_anchor = anchor_map.get(anchor, MSO_ANCHOR.MIDDLE)

        p = tf.paragraphs[0]
        align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
        p.alignment = align_map.get(align, PP_ALIGN.LEFT)

        spec = normalize_text_spec(text)
        effective_line_spacing = spec.get("line_spacing", spec.get("interligne", line_spacing))
        effective_line_spacing = normalize_line_spacing(effective_line_spacing)
        if effective_line_spacing is not None:
            p.line_spacing = effective_line_spacing

        base_style = {
            "font": font,
            "size_pt": size_pt,
            "bold": bold,
            "italic": italic,
            "color_hex": color_hex or text_color_hex,
            "cap": cap,
            "alpha_pct": alpha_pct,
        }
        for segment, emphasis in emphasis_spans(spec["text"], spec["emphasis"]):
            if not segment:
                continue
            run = p.add_run()
            run.text = segment
            set_run_props(run, **run_style(base_style, emphasis, emphasis_style))
    return shp


def add_oval(slide, x_cm, y_cm, w_cm, h_cm, fill_hex, *, line_hex=None, line_w_pt=0,
             name=None, text=None, font="Segoe UI", size_pt=12, bold=False,
             italic=False, color_hex="#070E1D", align="center", anchor="m"):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, cm(x_cm), cm(y_cm), cm(w_cm), cm(h_cm))
    remove_theme_style(shp)
    if name:
        shp.name = name
    shp.fill.solid()
    shp.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    if line_hex:
        shp.line.color.rgb = hex_to_rgb(line_hex)
        shp.line.width = Pt(line_w_pt) if line_w_pt else Pt(0.25)
    else:
        shp.line.fill.background()
    if text is not None:
        tf = shp.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)

        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        anchor_map = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}
        tf.vertical_anchor = anchor_map.get(anchor, MSO_ANCHOR.MIDDLE)

        p = tf.paragraphs[0]
        align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
        p.alignment = align_map.get(align, PP_ALIGN.CENTER)

        run = p.add_run()
        run.text = str(text)
        set_run_props(run, font=font, size_pt=size_pt, bold=bold, italic=italic,
                      color_hex=color_hex)
    return shp


def _load_pil_font(font, size_px, bold=False):
    windows_fonts = Path("C:/Windows/Fonts")
    font_key = (font or "").lower()
    candidates = []
    if "segoe ui black" in font_key:
        candidates = ["seguibl.ttf", "seguibl.TTF"]
    elif "segoe ui" in font_key:
        candidates = ["segoeuib.ttf" if bold else "segoeui.ttf", "seguibl.ttf"]
    elif "calibri" in font_key:
        candidates = ["calibrib.ttf" if bold else "calibri.ttf"]
    candidates.extend(["arialbd.ttf" if bold else "arial.ttf", "DejaVuSans-Bold.ttf" if bold else "DejaVuSans.ttf"])

    for candidate in candidates:
        try:
            return ImageFont.truetype(candidate, size_px)
        except OSError:
            pass
        font_path = windows_fonts / candidate
        if font_path.exists():
            try:
                return ImageFont.truetype(str(font_path), size_px)
            except OSError:
                pass
    return ImageFont.load_default()


def add_transparent_text_image(slide, x_cm, y_cm, w_cm, h_cm, text, *, font="Segoe UI",
                               size_pt=72, bold=False, color_hex="#070E1D",
                               alpha_pct=40, align="left", anchor="t", name=None):
    scale = 3
    px_w = max(1, int(w_cm / 2.54 * 96 * scale))
    px_h = max(1, int(h_cm / 2.54 * 96 * scale))
    img = Image.new("RGBA", (px_w, px_h), (255, 255, 255, 0))
    draw = ImageDraw.Draw(img)
    pil_font = _load_pil_font(font, int(size_pt / 72 * 96 * scale), bold=bold)
    bbox = draw.textbbox((0, 0), str(text), font=pil_font)
    text_w = bbox[2] - bbox[0]
    text_h = bbox[3] - bbox[1]

    align_map = {
        "left": 0,
        "center": (px_w - text_w) / 2,
        "right": px_w - text_w,
    }
    anchor_map = {
        "t": 0,
        "m": (px_h - text_h) / 2,
        "b": px_h - text_h,
    }
    x = align_map.get(align, 0) - bbox[0]
    y = anchor_map.get(anchor, 0) - bbox[1]
    draw.text((x, y), str(text), font=pil_font, fill=hex_to_rgba(color_hex, alpha_pct))

    buffer = BytesIO()
    img.save(buffer, format="PNG")
    buffer.seek(0)
    shp = slide.shapes.add_picture(buffer, cm(x_cm), cm(y_cm), width=cm(w_cm), height=cm(h_cm))
    remove_theme_style(shp)
    if name:
        shp.name = name
    return shp


def add_teardrop(slide, x_cm, y_cm, w_cm, h_cm, fill_hex, *, name=None):
    """Larme (teardrop). En python-pptx : MSO_SHAPE.TEAR (rotation à appliquer pour pointe en haut)."""
    shp = slide.shapes.add_shape(MSO_SHAPE.TEAR, cm(x_cm), cm(y_cm), cm(w_cm), cm(h_cm))
    remove_theme_style(shp)
    if name:
        shp.name = name
    # rotation pour pointe vers le bas (esthétique larme du template)
    shp.rotation = 180
    shp.fill.solid()
    shp.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    shp.line.fill.background()
    return shp


def add_line(slide, x1_cm, y1_cm, x2_cm, y2_cm, color_hex, *, width_pt=1.0, dash=False):
    from pptx.enum.shapes import MSO_CONNECTOR
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, cm(x1_cm), cm(y1_cm),
                                       cm(x2_cm), cm(y2_cm))
    remove_theme_style(conn)
    conn.line.color.rgb = hex_to_rgb(color_hex)
    conn.line.width = Pt(width_pt)
    if dash:
        # injection dash dans XML
        ln = conn.line._get_or_add_ln()
        prstDash = etree.SubElement(ln, qn("a:prstDash"))
        prstDash.set("val", "dash")
    return conn


# ---------- Bandeau Take-away avec gradient vertical 3-stops ----------

def add_takeaway_band(slide, text="Take away", *, y_cm=None):
    """Bandeau bordeaux avec gradient vertical 3-stops + texte centré blanc."""
    POS_X, POS_Y, W, H = 1.09, 16.77, 31.70, 1.10
    if y_cm is not None:
        POS_Y = y_cm
    BASE = "A25871"

    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cm(POS_X), cm(POS_Y), cm(W), cm(H))
    remove_theme_style(shp)
    shp.name = "TakeawayBand"

    # spPr est directement le CT_ShapeProperties accessible sur shp._element.spPr
    sppr = shp._element.spPr
    # nettoyer fill/line existants pour garder l'ordre OOXML valide :
    # géométrie -> fill -> line. Sinon PowerPoint peut ignorer le gradient
    # et retomber sur la couleur accent du thème.
    for tag in ["a:solidFill", "a:gradFill", "a:noFill", "a:blipFill", "a:pattFill", "a:ln"]:
        for el in sppr.findall(qn(tag)):
            sppr.remove(el)

    gradFill = etree.SubElement(sppr, qn("a:gradFill"))
    gradFill.set("flip", "none")
    gradFill.set("rotWithShape", "1")
    gsLst = etree.SubElement(gradFill, qn("a:gsLst"))

    # 3 stops : 0% (haut, plus clair), 50%, 100% (bas, plus sombre)
    stops = [(0, 30000), (50000, 67500), (100000, 100000)]
    for pos, shade in stops:
        gs = etree.SubElement(gsLst, qn("a:gs"))
        gs.set("pos", str(pos))
        srgb = etree.SubElement(gs, qn("a:srgbClr"))
        srgb.set("val", BASE)
        sm = etree.SubElement(srgb, qn("a:satMod"))
        sm.set("val", "115000")
        sh = etree.SubElement(srgb, qn("a:shade"))
        sh.set("val", str(shade))

    # angle vertical (5400000 = 90°)
    lin = etree.SubElement(gradFill, qn("a:lin"))
    lin.set("ang", "5400000")
    lin.set("scaled", "1")

    ln = etree.SubElement(sppr, qn("a:ln"))
    etree.SubElement(ln, qn("a:noFill"))

    # Texte centré
    add_textbox(slide, POS_X, POS_Y, W, H, text,
                font="Segoe UI", size_pt=16, bold=True, color_hex="#FFFFFF",
                align="center", anchor="m", name="TakeawayBandText")
    return shp


# ---------- Master commun (titre, signature bar, footer, logo, page num) ----------

def add_master_chrome(slide, title_text, *, page_num=None, total=None, with_signature=True,
                       logo_path=None):
    """Ajoute le 'chrome' commun à toutes les slides : titre, barre signature, footer, logo, n° page."""
    # Titre
    add_textbox(slide, 1.09, 0.71, 31.70, 1.58, title_text,
                font="Segoe UI Black", size_pt=29.32, color_hex="#484C6A",
                align="left", anchor="m", margin_left=0, name="SLIDE_TITLE")

    # Barre signature double
    if with_signature:
        # rectangle bleu nuit
        add_rect(slide, 1.09, 2.46, 1.87, 0.27, "#484C6A", name="SignatureRect")
        # trait long fin (simulé par un rectangle très plat)
        add_rect(slide, 1.09, 2.71, 3.69, 0.025, "#112753", name="SignatureLine")

    # Footer confidentiel
    footer_text = ("This document and the information therein are the property of Safran. "
                   "They must not be copied or communicated to a third party without the "
                   "prior written authorization of Safran")
    add_textbox(slide, 1.06, 18.51, 26.32, 0.40, footer_text,
                font="Segoe UI", size_pt=9, color_hex="#A8A8B5",
                align="left", anchor="t", name="FooterConfidential")

    # Numéro de page
    if page_num is not None:
        page_str = f"{page_num:02d}" if total is None else  f"{page_num:02d}{' ' * 6}|"
        add_textbox(slide, 1.04, 18.03, 1.50, 0.46, page_str,
                    font="Segoe UI", size_pt=9, color_hex="#7D7B8D",
                    align="left", anchor="m", name="PageNumber")

    # Logo Safran (si fourni)
    if logo_path:
        try:
            slide.shapes.add_picture(logo_path, cm(29.77), cm(17.91),
                                     width=cm(3.03), height=cm(0.64))
        except Exception:
            pass


# ---------- Fond blanc explicite ----------

def set_white_background(slide, color_hex="#FFFFFF"):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color_hex)
