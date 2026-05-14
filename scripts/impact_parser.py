#!/usr/bin/env python3
"""
Parseur de fichiers d'analyse d'impact PPTX (format Safran).

Extrait les données des zones nommées TITLE, ZONE_OMOC et RESSORTS_CHANGE
de chaque slide pour générer des entrées Layout 15 (Matrice d'impact par population).

Usage standalone :
    python impact_parser.py <chemin_analyse_impact.pptx> [--json]
"""
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

try:
    import fitz
    HAS_FITZ = True
except ImportError:
    HAS_FITZ = False

try:
    import pdfplumber
    HAS_PDFPLUMBER = True
except ImportError:
    HAS_PDFPLUMBER = False


# ---------------------------------------------------------------------------
# Constantes géométriques du radar OMOC (GAP MEASURE)
# ---------------------------------------------------------------------------
# Le radar est un diamant centré dans l'image ZONE_OMOC.
# Image : pos=(0.0, 3.25) cm, taille=(11.37, 9.57) cm
# Centre du diamant (en coordonnées slide) :
OMOC_CENTER_X = 5.685   # cm
OMOC_CENTER_Y = 8.035   # cm

# Pas unique (calibré empiriquement sur le template Safran).
# Distance en cm entre deux niveaux concentriques du radar.
OMOC_STEP = 0.93  # cm par niveau

# Noms des Ellipses portant les bullets OMOC (identiques sur chaque slide)
OMOC_BULLET_NAMES = {'5', '7', '10', '14'}

# Axes du radar :
#   Organisation = axe vertical vers le haut (y décroissant)
#   Métier       = axe horizontal vers la droite (x croissant)
#   Outils       = axe vertical vers le bas (y croissant)
#   Culture      = axe horizontal vers la gauche (x décroissant)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _classify_omoc_bullet(x_cm, y_cm, w_cm, h_cm):
    """
    Détermine l'axe OMOC et le niveau (0-4) d'un bullet à partir de sa
    position et taille sur le slide.

    Retourne (axis_name, level) où axis_name ∈ {Organization, Business, Tool, Culture}.
    """
    # Centre du bullet
    cx = x_cm + w_cm / 2
    cy = y_cm + h_cm / 2

    dx = cx - OMOC_CENTER_X
    dy = OMOC_CENTER_Y - cy  # inversé : y monte = dy positif

    # Déterminer l'axe le plus proche
    if abs(dx) > abs(dy):
        axis = 'Business' if dx > 0 else 'Culture'
        distance = abs(dx)
    else:
        axis = 'Organization' if dy > 0 else 'Tool'
        distance = abs(dy)

    level = round(distance / OMOC_STEP)
    level = max(0, min(4, level))
    return axis, level


def _extract_raw_levers(table):
    """
    Extrait les données brutes de la table RESSORTS_CHANGE.

    Retourne une liste de dicts :
    [
        {"ressort": "Sens", "type": "Communication", "detail": "...texte complet..."},
        ...
    ]
    """
    entries = []
    for ri in range(1, len(table.rows)):
        ressort = table.cell(ri, 0).text.strip()
        action_raw = table.cell(ri, 1).text.strip()
        if not ressort or not action_raw:
            continue

        lines = action_raw.split('\n')
        action_type = lines[0].strip().rstrip(':')
        detail = ' '.join(l.strip() for l in lines[1:] if l.strip())
        entries.append({
            "ressort": ressort,
            "type": action_type,
            "detail": detail or action_type
        })
    return entries


def _fallback_summarize_levers(raw_levers, max_len=135):
    """
    Résumé mécanique de secours (≤ max_len caractères) utilisé quand
    l'IA ne fournit pas de résumé rédigé.

    Concatène les types d'action uniques, puis complète avec les noms
    de ressorts si la place le permet.
    """
    if not raw_levers:
        return ""

    # Types d'action uniques (Communication, Support, Mobilisation, Formation)
    types_seen = []
    types_set = set()
    for entry in raw_levers:
        t = entry["type"]
        if t.lower() not in types_set:
            types_set.add(t.lower())
            types_seen.append(t)

    # Résumé court : "Type1 | Type2 | ..."
    summary = " | ".join(types_seen)
    if len(summary) <= max_len:
        return summary

    # Tronquer si trop long
    return summary[:max_len - 1].rsplit(' ', 1)[0] + '…'


def _parse_title(title_text):
    """
    Sépare le texte du TITLE en (population, effectif).

    Formats attendus :
      - "Population – N personnes"
      - "Population - N personnes"
      - "Population" (sans effectif)

    Retourne (population, effectif) où effectif = "TBD" si absent.
    """
    # Essayer en-dash d'abord, puis tiret simple
    for sep in ['\u2013', '-']:
        if sep in title_text:
            parts = title_text.split(sep, 1)
            population = parts[0].strip()
            effectif_raw = parts[1].strip()
            if effectif_raw:
                return population, f"~ {effectif_raw}"
            break
    else:
        population = title_text.strip()

    return population, "TBD"


# ---------------------------------------------------------------------------
# Parsing principal
# ---------------------------------------------------------------------------

def parse_impact_file(pptx_path):
    """
    Parse un fichier d'analyse d'impact PPTX et retourne une liste de
    dictionnaires, un par slide/population :

    [
        {
            "population": "Garants de fabricabilité",
            "effectif": "~ 5 personnes",
            "omoc": {"Tool": 1, "Business": 1, "Organization": 2, "Culture": 2},
            "raw_levers": [{"ressort": "Sens", "type": "Communication", "detail": "..."}],
            "levers": "Communication | Support"   # fallback mécanique
        },
        ...
    ]
    """
    prs = Presentation(str(pptx_path))
    populations = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        entry = {
            "slide_index": slide_idx,
            "population": "",
            "effectif": "TBD",
            "omoc": {"Tool": 0, "Business": 0, "Organization": 0, "Culture": 0},
            "raw_levers": [],
            "levers": ""
        }

        for shape in slide.shapes:
            name = shape.name.strip()

            # --- TITLE / TITRE ---
            if name in ('TITLE', 'TITRE'):
                title_text = shape.text_frame.text.strip()
                entry["population"], entry["effectif"] = _parse_title(title_text)

            # --- ZONE_OMOC bullets (Ellipse 5, 7, 10, 14) ---
            if 'Ellipse' in name:
                parts = name.split()
                if len(parts) >= 2 and parts[-1] in OMOC_BULLET_NAMES:
                    x = Emu(shape.left).cm
                    y = Emu(shape.top).cm
                    w = Emu(shape.width).cm
                    h = Emu(shape.height).cm
                    axis, level = _classify_omoc_bullet(x, y, w, h)
                    entry["omoc"][axis] = level

            # --- RESSORTS_CHANGE ---
            if name == 'RESSORTS_CHANGE':
                entry["raw_levers"] = _extract_raw_levers(shape.table)
                entry["levers"] = _fallback_summarize_levers(entry["raw_levers"])

        # N'ajouter que les slides qui ont un titre exploitable
        if entry["population"]:
            populations.append(entry)

    return populations


def _parse_impact_pdf_fitz(pdf_path):
    """Parse un PDF d'analyse d'impact via PyMuPDF (fitz)."""
    pdf = fitz.open(str(pdf_path))
    populations = []

    def is_red(fill):
        if fill is None: return False
        return fill[0] > 0.9 and fill[1] < 0.1 and fill[2] < 0.1

    def is_title_block(x0, y0, x1, y1, text, page_w, page_h):
        return (0.02 * page_w <= x0 <= 0.60 * page_w and
                0.02 * page_h <= y0 <= 0.15 * page_h and
                len(text.strip()) > 5)

    for slide_idx in range(len(pdf)):
        page = pdf[slide_idx]
        entry = {
            "slide_index": slide_idx + 1,
            "population": "",
            "effectif": "TBD",
            "omoc": {"Tool": 0, "Business": 0, "Organization": 0, "Culture": 0},
            "raw_levers": [],
            "levers": ""
        }

        page_rect = page.rect
        page_w = page_rect.width
        page_h = page_rect.height

        blocks = page.get_text("blocks")
        title = ""
        for b in blocks:
            x0, y0, x1, y1, text = b[0], b[1], b[2], b[3], b[4].strip()
            if is_title_block(x0, y0, x1, y1, text, page_w, page_h):
                title = text.replace('\n', ' ')
                break

        if title:
            entry["population"], entry["effectif"] = _parse_title(title)
        else:
            continue

        drawings = page.get_drawings()
        for d in drawings:
            fill = d.get("fill")
            r = d.get("rect")
            if r and is_red(fill):
                cx = (r.x0 + r.x1) / 2
                cy = (r.y0 + r.y1) / 2
                
                px = (cx / page_w) * 100
                py = (cy / page_h) * 100
                
                if 5 <= px <= 40 and 20 <= py <= 65:
                    dx = px - 16.8
                    dy = 42.2 - py
                    
                    level_x = abs(dx) / 2.75
                    level_y = abs(dy) / 4.88
                    
                    if abs(dx) < 1.0 and abs(dy) < 1.0:
                        continue
                        
                    if level_x > level_y:
                        axis = 'Business' if dx > 0 else 'Culture'
                        level = round(level_x)
                    else:
                        axis = 'Organization' if dy > 0 else 'Tool'
                        level = round(level_y)
                        
                    entry["omoc"][axis] = max(0, min(4, level))

        for b in blocks:
            x0, y0, x1, y1, text = b[0], b[1], b[2], b[3], b[4].strip()
            if 0.35 * page_w <= x0 <= 0.95 * page_w and 0.50 * page_h <= y0 <= 0.95 * page_h:
                lines = [l.strip() for l in text.split('\n') if l.strip()]
                if len(lines) >= 3:
                    ressort = lines[0]
                    if ressort.lower() in ('ressorts du changement', 'actions'):
                        continue
                    
                    action_type = lines[1]
                    if lines[2] in ['•', '-', '', '\uf0b7', '']:
                        detail = " ".join(lines[3:])
                    else:
                        detail = " ".join(lines[2:]).lstrip('•-\uf0b7 ')
                    
                    entry["raw_levers"].append({
                        "ressort": ressort,
                        "type": action_type.rstrip(':'),
                        "detail": detail or action_type
                    })
        
        entry["levers"] = _fallback_summarize_levers(entry["raw_levers"])
        populations.append(entry)

    pdf.close()
    return populations


def _parse_impact_pdf_pdfplumber(pdf_path):
    """Parse un PDF d'analyse d'impact via pdfplumber (alternative à PyMuPDF)."""

    def _group_into_lines(words, y_tol=3):
        words = sorted(words, key=lambda w: (w["top"], w["x0"]))
        lines = []
        for w in words:
            if not lines or w["top"] - lines[-1][0]["top"] > y_tol:
                lines.append([w])
            else:
                lines[-1].append(w)
        return [
            (" ".join(w["text"] for w in line),
             min(w["top"] for w in line),
             max(w["bottom"] for w in line))
            for line in lines
        ]

    def _lines_to_blocks(lines, gap=12):
        blocks, current = [], [lines[0]]
        for line in lines[1:]:
            if line[1] - current[-1][2] > gap:
                blocks.append(current)
                current = [line]
            else:
                current.append(line)
        blocks.append(current)
        return ["\n".join(l[0] for l in block) for block in blocks]

    populations = []
    with pdfplumber.open(str(pdf_path)) as pdf:
        for slide_idx, page in enumerate(pdf.pages):
            entry = {
                "slide_index": slide_idx + 1,
                "population": "",
                "effectif": "TBD",
                "omoc": {"Tool": 0, "Business": 0, "Organization": 0, "Culture": 0},
                "raw_levers": [],
                "levers": ""
            }

            page_w = page.width
            page_h = page.height

            words = page.extract_words(keep_blank_chars=False, x_tolerance=3, y_tolerance=3)

            # --- Titre (Zone 1) ---
            title_words = [
                w for w in words 
                if 0.02 * page_w <= w["x0"] <= 0.60 * page_w 
                and 0.02 * page_h <= w["top"] <= 0.15 * page_h
            ]
            if not title_words:
                continue
            title = " ".join(
                w["text"] for w in sorted(title_words, key=lambda w: (w["top"], w["x0"]))
            ).replace('\n', ' ')
            entry["population"], entry["effectif"] = _parse_title(title)

            # --- Bullets OMOC (Zone 2) ---
            for rect in page.rects:
                cx = (rect["x0"] + rect["x1"]) / 2
                cy = (rect["top"] + rect["bottom"]) / 2
                
                px = (cx / page_w) * 100
                py = (cy / page_h) * 100
                
                if not (5 <= px <= 40 and 20 <= py <= 65):
                    continue

                fill = rect.get("non_stroking_color")
                if fill is not None:
                    if isinstance(fill, (list, tuple)) and len(fill) >= 3:
                        r, g, b = fill[0], fill[1], fill[2]
                        if not (r > 0.9 and g < 0.1 and b < 0.1):
                            continue
                    else:
                        continue

                dx = px - 16.8
                dy = 42.2 - py
                
                level_x = abs(dx) / 2.75
                level_y = abs(dy) / 4.88
                
                if abs(dx) < 1.0 and abs(dy) < 1.0:
                    continue

                if level_x > level_y:
                    axis = 'Business' if dx > 0 else 'Culture'
                    level = round(level_x)
                else:
                    axis = 'Organization' if dy > 0 else 'Tool'
                    level = round(level_y)
                entry["omoc"][axis] = max(0, min(4, level))

            # --- Leviers (Zone 3) ---
            lever_words = [
                w for w in words 
                if 0.35 * page_w <= w["x0"] <= 0.95 * page_w 
                and 0.50 * page_h <= w["top"] <= 0.95 * page_h
            ]
            if lever_words:
                lines = _group_into_lines(lever_words)
                for block_text in _lines_to_blocks(lines):
                    block_lines = [l.strip() for l in block_text.split('\n') if l.strip()]
                    if len(block_lines) < 3:
                        continue
                    ressort = block_lines[0]
                    if ressort.lower() in ('ressorts du changement', 'actions'):
                        continue
                    action_type = block_lines[1]
                    if block_lines[2] in ['•', '-', '', '', '\uf0b7']:
                        detail = " ".join(block_lines[3:])
                    else:
                        detail = " ".join(block_lines[2:]).lstrip('•-\uf0b7 ')
                    entry["raw_levers"].append({
                        "ressort": ressort,
                        "type": action_type.rstrip(':'),
                        "detail": detail or action_type
                    })

            entry["levers"] = _fallback_summarize_levers(entry["raw_levers"])
            populations.append(entry)

    return populations


def parse_impact_pdf(pdf_path):
    """
    Parse un fichier d'analyse d'impact PDF et retourne la même structure
    que parse_impact_file(). Utilise PyMuPDF si disponible, sinon pdfplumber.
    """
    if HAS_FITZ:
        return _parse_impact_pdf_fitz(pdf_path)
    if HAS_PDFPLUMBER:
        return _parse_impact_pdf_pdfplumber(pdf_path)
    print(
        "[ERROR] PyMuPDF (fitz) ou pdfplumber est requis pour parser les PDF d'analyse d'impact.\n"
        "        Installez l'un des deux : pip install pymupdf  ou  pip install pdfplumber",
        file=sys.stderr
    )
    return []


# ---------------------------------------------------------------------------
# Génération du contenu Layout 15
# ---------------------------------------------------------------------------

def build_layout15_slides(populations):
    """
    Convertit la liste de populations en entrées content.json pour le Layout 15.

    Règles :
    - Maximum 6 populations par slide Layout 15
    - Si > 6 populations, créer des slides supplémentaires
    - Mapping :
        label    → nom de la population (S15_TABLE_LABEL_Ri)
        sublabel → effectif (S15_TABLE_SUBLABEL_Ri)
        impacts  → [Outils, Métier, Organisation, Culture] (S15_TABLE_Ri_Ci)
        levers   → résumé RESSORTS_CHANGE (S15_TABLE_Ri_C6)
    """
    slides = []

    # Chunking par paquets de 6
    for chunk_start in range(0, len(populations), 6):
        chunk = populations[chunk_start:chunk_start + 6]

        rows = []
        for pop in chunk:
            omoc = pop["omoc"]
            rows.append({
                "label": pop["population"],
                "sublabel": pop["effectif"],
                "impacts": [
                    omoc.get("Tool", 0),         # Outils
                    omoc.get("Business", 0),      # Métier
                    omoc.get("Organization", 0),  # Organisation
                    omoc.get("Culture", 0),       # Culture
                ],
                "levers": pop["levers"]
            })

        slide_entry = {
            "layout": 15,
            "content": {
                "title": "Impact par population",
                "headers": [
                    "Outils", "Métier", "Organisation", "Culture",
                    "Leviers retenus"
                ],
                "rows": rows,
                "legend_title": "Intensité",
                "legend": ["Faible", "Modéré", "Fort", "Très fort"],
                "footnote": "Leviers = axes d'accompagnement prioritaires par population"
            }
        }
        slides.append(slide_entry)

    return slides


def generate_impact_content(file_path):
    """
    Fonction principale : parse le fichier d'analyse d'impact et retourne
    la liste de slides Layout 15 prête à être insérée dans le content.json.
    """
    path = Path(file_path)
    if path.suffix.lower() == ".pdf":
        populations = parse_impact_pdf(file_path)
    else:
        populations = parse_impact_file(file_path)
    return build_layout15_slides(populations)


# ---------------------------------------------------------------------------
# Mode standalone
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python impact_parser.py <fichier_analyse_impact> [--json]",
              file=sys.stderr)
        sys.exit(1)

    file_path = sys.argv[1]
    output_json = "--json" in sys.argv

    path = Path(file_path)
    if path.suffix.lower() == ".pdf":
        populations = parse_impact_pdf(file_path)
    else:
        populations = parse_impact_file(file_path)

    if output_json:
        slides = build_layout15_slides(populations)
        print(json.dumps({"slides": slides}, ensure_ascii=False, indent=2))
    elif "--raw" in sys.argv:
        # Mode brut : afficher les données RESSORTS_CHANGE complètes (pour l'IA)
        print(json.dumps(populations, ensure_ascii=False, indent=2))
    else:
        print(f"Populations extraites : {len(populations)}")
        print("=" * 70)
        for i, pop in enumerate(populations, 1):
            omoc = pop["omoc"]
            si = pop.get("slide_index", "?")
            print(f"\n[{i}] (slide source #{si}) {pop['population']}")
            print(f"    Effectif     : {pop['effectif']}")
            print(f"    Outils       : {omoc['Tool']}")
            print(f"    Metier       : {omoc['Business']}")
            print(f"    Organisation : {omoc['Organization']}")
            print(f"    Culture      : {omoc['Culture']}")
            print(f"    Leviers (fallback) : {pop['levers']}")
            if pop['raw_levers']:
                print(f"    Actions brutes :")
                for rl in pop['raw_levers']:
                    print(f"      - [{rl['ressort']}] {rl['type']}: {rl['detail'][:80]}")

        print(f"\n=> {len(populations)} populations => "
              f"{(len(populations) + 5) // 6} slide(s) Layout 15")
